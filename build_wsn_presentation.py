from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


OUT_PPTX = "sensornet-project-wsn-v3-dhcpv6-presentation.pptx"
OUT_NOTES = "sensornet-project-wsn-v3-dhcpv6-presentation-notes.md"

NAVY = RGBColor(20, 45, 75)
BLUE = RGBColor(42, 104, 163)
TEAL = RGBColor(24, 132, 127)
GREEN = RGBColor(69, 150, 86)
ORANGE = RGBColor(221, 126, 49)
RED = RGBColor(190, 69, 69)
GRAY = RGBColor(88, 96, 105)
LIGHT = RGBColor(244, 247, 250)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(20, 24, 28)


def set_font(run, size=20, bold=False, color=BLACK):
    run.font.name = "Noto Sans CJK KR"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.28), Inches(12.45), Inches(0.72))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.LEFT
    set_font(p.runs[0], 27, True, NAVY)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.48), Inches(0.95), Inches(12.2), Inches(0.36))
        p = sub.text_frame.paragraphs[0]
        p.text = subtitle
        set_font(p.runs[0], 12, False, GRAY)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.45), Inches(1.22), Inches(12.45), Inches(0.035))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()


def add_footer(slide, num):
    box = slide.shapes.add_textbox(Inches(11.9), Inches(7.05), Inches(1.0), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    p.text = str(num)
    p.alignment = PP_ALIGN.RIGHT
    set_font(p.runs[0], 9, False, GRAY)


def add_bullets(slide, x, y, w, h, bullets, size=17, color=BLACK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(7)
        set_font(p.runs[0], size, False, color)
    return box


def add_card(slide, x, y, w, h, title, body, fill=LIGHT, accent=BLUE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RGBColor(218, 225, 232)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.09), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.12), Inches(w - 0.32), Inches(0.32))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    set_font(p.runs[0], 15, True, NAVY)
    add_bullets(slide, x + 0.18, y + 0.52, w - 0.34, h - 0.65, body, size=12)


def add_label(slide, x, y, w, h, text, fill, size=13, color=WHITE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    set_font(p.runs[0], size, True, color)
    return shape


def add_arrow(slide, x, y, w, h, color=GRAY):
    arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()
    return arrow


def add_code(slide, x, y, w, h, lines):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(31, 38, 46)
    shape.line.color.rgb = RGBColor(31, 38, 46)
    tb = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.14), Inches(w - 0.32), Inches(h - 0.22))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line if line else " "
        p.space_after = Pt(3)
        run.font.name = "DejaVu Sans Mono"
        run.font.size = Pt(10)
        run.font.color.rgb = WHITE


def add_two_column(slide, left_title, left_bullets, right_title, right_bullets):
    add_card(slide, 0.75, 1.65, 5.85, 4.75, left_title, left_bullets, accent=TEAL)
    add_card(slide, 6.85, 1.65, 5.85, 4.75, right_title, right_bullets, accent=ORANGE)


slides = [
    {
        "title": "DHCPv6 기반 IPv6/RPL Wireless Sensor Network",
        "subtitle": "TinyOS + TelosB/Kmote | PPP Border Router + DHCPv6 Address Assignment + Sensor Nodes",
        "type": "cover",
    },
    {
        "title": "프로젝트 요구사항과 구현 범위",
        "subtitle": "sensornet-project-wsn-v3.pdf 기준",
        "type": "requirements",
    },
    {
        "title": "설계 목표: 자동 주소 할당이 가능한 WSN",
        "subtitle": "강의의 IP narrow waist와 DHCPv6/RPL 구조를 WSN 프로젝트에 적용",
        "type": "goals",
    },
    {
        "title": "전체 시스템 구조",
        "subtitle": "PC DHCPv6 server -> PPP border router -> IEEE 802.15.4 sensor nodes",
        "type": "architecture",
    },
    {
        "title": "프로토콜 스택과 강의 내용 연결",
        "subtitle": "IP, DHCPv6, 6LoWPAN, RPL, TCP/UDP를 TinyOS/BLIP 위에 구성",
        "type": "stack",
    },
    {
        "title": "IPv6 주소가 없을 때 DHCPv6는 어떻게 시작되는가",
        "subtitle": "Active Message가 아니라 IPv6 link-local + UDP multicast로 bootstrap",
        "type": "ipv6_bootstrap",
    },
    {
        "title": "DHCPv6 주소 할당 메시지 흐름",
        "subtitle": "sensor node client -> node1 relay -> PC dnsmasq server",
        "type": "dhcp_flow",
    },
    {
        "title": "RPL 경로와 TCP/HTTP 동작",
        "subtitle": "주소 할당 이후에는 일반 IPv6 packet처럼 RPL route를 따라 전달",
        "type": "rpl_tcp_flow",
    },
    {
        "title": "Part-1: Command 지원",
        "subtitle": "프로젝트 필수 명령을 HTTP API로 매핑",
        "type": "commands",
    },
    {
        "title": "센서 데이터 수집 경로",
        "subtitle": "주기적 report와 요청 기반 query를 동시에 지원",
        "type": "sensor",
    },
    {
        "title": "Part-2: Advanced Feature",
        "subtitle": "DHCPv6 + RPL + TCP + Internet-facing 접근 구조",
        "type": "advanced",
    },
    {
        "title": "구현 상세",
        "subtitle": "TinyOS 컴포넌트와 빌드 설정",
        "type": "implementation",
    },
    {
        "title": "빌드 및 설치 결과",
        "subtitle": "현재 연결된 TelosB/Kmote 장치 기준 업데이트 완료",
        "type": "build",
    },
    {
        "title": "데모 시나리오",
        "subtitle": "발표 영상에 넣을 순서와 명령어",
        "type": "demo",
    },
    {
        "title": "한계와 개선 방향",
        "subtitle": "평가/확장 포인트를 강의 주제와 연결",
        "type": "future",
    },
    {
        "title": "결론",
        "subtitle": "프로젝트 요구사항 대응 요약",
        "type": "conclusion",
    },
]


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

for idx, spec in enumerate(slides, start=1):
    slide = prs.slides.add_slide(blank)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE
    add_title(slide, spec["title"], spec.get("subtitle"))
    add_footer(slide, idx)

    t = spec["type"]
    if t == "cover":
        add_label(slide, 0.8, 1.75, 4.2, 0.6, "Final Project", BLUE, 18)
        title = slide.shapes.add_textbox(Inches(0.8), Inches(2.55), Inches(8.7), Inches(1.05))
        p = title.text_frame.paragraphs[0]
        p.text = "Implementation of a Wireless Sensor Network"
        set_font(p.runs[0], 34, True, NAVY)
        add_bullets(slide, 0.85, 3.75, 6.9, 1.35, [
            "Gateway: PppRouter node 1, DHCPv6 relay + RPL DODAG root",
            "Sensor nodes: TCPEcho node 2/3/4",
            "DHCPv6-assigned IPv6 addresses, 6LoWPAN compression, RPL storing mode",
        ], 17)
        add_card(slide, 8.35, 2.05, 4.1, 3.8, "핵심 메시지", [
            "HW2 수준의 단순 수집 네트워크를 명령 기반 WSN으로 확장",
            "강의의 IP/DHCPv6/RPL/Transport 개념을 TinyOS mote에 통합",
            "DHCP lease로 받은 IPv6 주소에 TCP/HTTP command 전송",
        ], accent=TEAL)

    elif t == "requirements":
        add_two_column(slide, "Project PDF 요구사항", [
            "5 WSN motes 권장: 1 base station + sensor nodes",
            "Part-1: get-led, set-led, get-voltage, get-temp, get-channel, get-tx-power",
            "Part-2: multihop, RPL, TCP, Internet 연결 등 advanced feature 선택",
            "제출물: pptx report, presentation+demo mp4, GitHub link",
        ], "이번 구현 범위", [
            "현재 실험 장비: 4 motes 사용 가능",
            "node 1: PPP/RPL border router + DHCPv6 relay",
            "node 2/3/4: command + sensor + echo nodes",
            "advanced: DHCPv6/IPv6/6LoWPAN/RPL/TCP/HTTP 기반 원격 접근",
        ])

    elif t == "goals":
        add_card(slide, 0.75, 1.55, 3.9, 4.85, "Motivation", [
            "HW2는 base station에 센서 값이 일방향으로 모임",
            "최종 프로젝트는 사용자가 네트워크 내부 node에 명령을 보내야 함",
            "따라서 host와 WSN 사이의 양방향 접근성이 필요",
        ], accent=BLUE)
        add_card(slide, 4.95, 1.55, 3.9, 4.85, "Design Goal", [
            "각 mote가 DHCPv6로 IPv6 주소를 자동 획득",
            "HTTP/TCP/UDP command interface 제공",
            "RPL로 multi-hop 확장 가능한 라우팅 기반 확보",
        ], accent=TEAL)
        add_card(slide, 9.15, 1.55, 3.45, 4.85, "Lecture Link", [
            "IP narrow waist: application과 link technology 분리",
            "DHCPv6: server/relay/client 기반 주소 자동화",
            "RPL: low-power lossy network용 IPv6 routing",
        ], accent=ORANGE)

    elif t == "architecture":
        add_label(slide, 0.65, 2.25, 2.15, 0.75, "PC / DHCPv6 Server\nfd00:23:42:1::100", BLUE, 11)
        add_arrow(slide, 2.9, 2.45, 0.8, 0.32, GRAY)
        add_label(slide, 3.85, 2.05, 2.3, 1.1, "node 1\nPppRouter\nDHCPv6 Relay\nRPL Root", TEAL, 11)
        add_arrow(slide, 6.25, 2.45, 0.8, 0.32, GRAY)
        add_label(slide, 7.2, 1.45, 1.65, 0.8, "node 2\nTCPEcho", GREEN, 12)
        add_label(slide, 9.05, 2.55, 1.65, 0.8, "node 3\nTCPEcho", GREEN, 12)
        add_label(slide, 10.9, 1.45, 1.65, 0.8, "node 4\nTCPEcho", GREEN, 12)
        add_label(slide, 7.55, 4.25, 4.3, 0.65, "IEEE 802.15.4 / channel 26 / PAN group 0xCA", ORANGE, 12)
        add_bullets(slide, 0.85, 4.25, 5.55, 1.65, [
            "PPP link bridges host network and WSN",
            "PppRouter forwards IPv6 packets between PPP and radio",
            "PppRouter relays DHCPv6 traffic between WSN and PC",
            "RPL DODAG root receives DAO routes from sensor nodes",
        ], 15)
        add_bullets(slide, 7.15, 5.15, 5.55, 1.15, [
            "Sensor node IPv6 addresses are assigned from DHCPv6 lease pool",
            "HTTP control: tcp/80, TCP echo: tcp/7, UDP echo: udp/7",
        ], 15)

    elif t == "stack":
        layers = [
            ("Application", "HTTP commands, TCP/UDP echo, sensor report", BLUE),
            ("Transport", "TCP port 80/7, UDP port 7/7777, DHCPv6 UDP 546/547", TEAL),
            ("Routing", "RPL storing mode, OF0, DIO/DAO", GREEN),
            ("Network", "IPv6, stateful DHCPv6 address assignment", ORANGE),
            ("Adaptation", "6LoWPAN header compression HC_VERSION=6", RED),
            ("Link/PHY", "IEEE 802.15.4 CC2420 channel 26", GRAY),
        ]
        y = 1.55
        for name, desc, color in layers:
            add_label(slide, 1.0, y, 2.35, 0.55, name, color, 13)
            add_bullets(slide, 3.65, y + 0.04, 8.3, 0.45, [desc], 14)
            y += 0.78
        add_card(slide, 1.0, 6.15, 11.3, 0.72, "강의 내용 연결", [
            "IP lecture: IPv6 over WSN, DHCPv6, RFC 4944 6LoWPAN adaptation/compression, border router model",
            "Routing lecture: RPL DODAG root, DIO/DAO, Objective Function",
            "Transport lecture: constrained WSN에서 TCP/UDP와 congestion/reliability tradeoff",
        ], accent=BLUE)

    elif t == "ipv6_bootstrap":
        add_card(slide, 0.75, 1.55, 3.8, 4.95, "초기 상태", [
            "sensor node는 아직 global IPv6 주소가 없음",
            "하지만 IEEE 802.15.4 EUI-64 기반 link-local 주소는 사용 가능",
            "BLIP IPv6 stack이 UDP socket을 열고 DHCPv6 client를 시작",
            "따라서 TinyOS Active Message로 주소를 나눠주는 구조가 아님",
        ], accent=BLUE)
        add_card(slide, 4.85, 1.55, 3.8, 4.95, "DHCPv6 시작", [
            "Dhcp6ClientP가 UDP 546번 포트에 bind",
            "SOLICIT을 ff02::1:2 all DHCP relay/server multicast로 전송",
            "client id는 mote의 EUI-64 기반 DUID-LL 사용",
            "transaction id로 ADVERTISE/REPLY가 자기 요청인지 확인",
        ], accent=TEAL)
        add_card(slide, 8.95, 1.55, 3.55, 4.95, "주소 적용", [
            "REPLY의 IA_NA/IAADDR option에서 IPv6 주소 획득",
            "SetIPAddress.setAddress()로 BLIP global address 등록",
            "주소 valid lifetime이 끝나면 removeAddress 후 재시도",
            "주소가 생긴 뒤 HTTP/TCP/UDP application traffic 가능",
        ], accent=ORANGE)

    elif t == "dhcp_flow":
        add_label(slide, 0.65, 1.75, 2.0, 0.58, "node2~4\nDHCPv6 Client", GREEN, 11)
        add_arrow(slide, 2.85, 1.88, 0.75, 0.28, GRAY)
        add_label(slide, 3.8, 1.75, 2.05, 0.58, "node1\nDHCPv6 Relay", TEAL, 11)
        add_arrow(slide, 6.05, 1.88, 0.75, 0.28, GRAY)
        add_label(slide, 7.0, 1.75, 2.35, 0.58, "PC\nDHCPv6 Server", BLUE, 11)
        add_arrow(slide, 9.55, 1.88, 0.75, 0.28, GRAY)
        add_label(slide, 10.5, 1.75, 2.1, 0.58, "lease\nfd00:.../64", ORANGE, 11)
        add_card(slide, 0.75, 2.75, 5.85, 3.55, "Forward 방향", [
            "1. client: SOLICIT/REQUEST, UDP 546 -> 547",
            "2. 목적지: ff02::1:2 또는 relay/server",
            "3. node1 relay: client packet을 RELAY_FORW로 감싸서 upstream 전송",
            "4. PC dnsmasq: fd00:23:42:1::1000~1fff lease pool에서 주소 선택",
        ], accent=TEAL)
        add_card(slide, 6.85, 2.75, 5.85, 3.55, "Return 방향", [
            "5. server: RELAY_REPLY 안에 ADVERTISE/REPLY 포함",
            "6. node1 relay: peer_addr을 보고 원래 sensor node로 전달",
            "7. client: REPLY의 IA_NA/IAADDR을 파싱",
            "8. client: global IPv6 address 등록 후 RPL/TCP traffic 가능",
        ], accent=ORANGE)

    elif t == "rpl_tcp_flow":
        add_card(slide, 0.75, 1.55, 3.8, 4.95, "RPL 제어 평면", [
            "node1 PppRouter가 RootControl.setRoot() 호출",
            "RPL root가 DIO를 전파하며 DODAG 형성",
            "sensor node는 parent를 선택하고 default route를 구성",
            "storing mode에서는 DAO로 downward route가 root 방향에 설치",
        ], accent=GREEN)
        add_card(slide, 4.85, 1.55, 3.8, 4.95, "TCP 연결", [
            "TCP 자체는 Active Message가 아니라 IPv6 위의 transport protocol",
            "PC curl -> node IPv6:80 으로 SYN 전송",
            "node1은 PPP에서 받은 IPv6 packet을 radio/RPL 쪽으로 forwarding",
            "sensor node의 TcpSocketC가 SYN을 accept하고 HTTP 요청 처리",
        ], accent=BLUE)
        add_card(slide, 8.95, 1.55, 3.55, 4.95, "왜 RPL이 필요한가", [
            "1-hop이면 node1이 직접 전달 가능",
            "multi-hop이면 중간 sensor node가 parent/child route를 따라 전달",
            "RPL route가 없으면 DHCP lease가 있어도 IP packet 도달 불가",
            "따라서 DHCPv6는 주소, RPL은 경로, TCP는 신뢰성 있는 command channel",
        ], accent=ORANGE)

    elif t == "commands":
        add_card(slide, 0.75, 1.55, 5.85, 4.85, "필수 명령", [
            "GET /get-led -> 3 LED 상태를 0~7 decimal 값으로 반환",
            "GET /set-led/<n> -> n의 bit pattern으로 LED 제어",
            "GET /get-voltage -> battery voltage mV",
            "GET /get-temp -> internal temperature",
            "GET /get-channel -> radio channel",
            "GET /get-tx-power -> radio TX power",
        ], accent=TEAL)
        add_card(slide, 6.85, 1.55, 5.85, 4.85, "추가 명령", [
            "GET /get-humidity -> humidity raw value",
            "GET /get-sensor -> node, seq, time, temp, humidity, voltage tuple",
            "GET /get-time -> local runtime",
            "GET /get-stats -> HTTP/TCP/sensor/report counters",
            "DHCPv6 lease/log -> node별 실제 IPv6 주소 확인",
        ], accent=ORANGE)

    elif t == "sensor":
        add_label(slide, 0.8, 1.75, 2.2, 0.72, "SensorTimer\n10 sec", BLUE, 12)
        add_arrow(slide, 3.15, 1.95, 0.7, 0.3, GRAY)
        add_label(slide, 4.0, 1.75, 2.1, 0.72, "Temp/Hum/Batt\nRead chain", TEAL, 12)
        add_arrow(slide, 6.25, 1.95, 0.7, 0.3, GRAY)
        add_label(slide, 7.1, 1.75, 2.15, 0.72, "UDP report\nhost:7777", GREEN, 12)
        add_arrow(slide, 9.4, 1.95, 0.7, 0.3, GRAY)
        add_label(slide, 10.25, 1.75, 2.1, 0.72, "PC log / demo\ncollector", ORANGE, 12)
        add_card(slide, 1.0, 3.3, 5.45, 2.35, "주기적 수집", [
            "node별 sensor_seq 증가",
            "temperature 변환값과 raw ADC 값을 함께 제공",
            "battery는 mV와 raw value를 같이 표시",
        ], accent=BLUE)
        add_card(slide, 6.9, 3.3, 5.35, 2.35, "요청 기반 조회", [
            "HTTP command가 들어오면 최신 cached sensor state를 반환",
            "아직 sensor sample이 없으면 -1로 명확히 표시",
            "데모에서 curl로 node별 응답 확인 가능",
        ], accent=TEAL)

    elif t == "advanced":
        add_card(slide, 0.75, 1.55, 3.8, 4.9, "DHCPv6 Addressing", [
            "IN6_PREFIX 제거로 DHCPv6 mode 사용",
            "PC dnsmasq가 fd00:23:42:1::/64 pool 제공",
            "PppRouter가 DHCPv6 client/relay로 동작",
            "sensor nodes는 DHCPv6 lease로 global IPv6 획득",
        ], accent=GREEN)
        add_card(slide, 4.85, 1.55, 3.8, 4.9, "RPL Routing", [
            "node 1이 RootControl.setRoot() 호출",
            "sensor nodes는 RPL_ROUTING=1",
            "storing mode + OF0로 hop-count 기반 parent 선택",
            "DHCP 이후 RPL route로 command path 형성",
        ], accent=BLUE)
        add_card(slide, 8.95, 1.55, 3.55, 4.9, "TCP Support", [
            "각 sensor node가 TCP echo port 7 제공",
            "HTTP command interface는 TCP port 80 사용",
            "Transport lecture의 TCP/UDP tradeoff를 실제 mote에서 확인",
        ], accent=ORANGE)

    elif t == "implementation":
        add_card(slide, 0.75, 1.55, 5.85, 4.95, "PppRouter", [
            "PppDaemonC + PppIpv6C로 PPP IPv6 link 구성",
            "IPForwardingEngineP와 연결하여 PPP/radio forwarding",
            "RPLRoutingC를 root로 설정",
            "Dhcp6C/Dhcp6ClientC로 DHCPv6 client/relay 구성",
        ], accent=BLUE)
        add_card(slide, 6.85, 1.55, 5.85, 4.95, "TCPEcho", [
            "UdpSocketC: echo/status/sensor report",
            "TcpSocketC: TCP echo와 HTTP command server",
            "SensirionSht11C + VoltageC sensor read",
            "Dhcp6C로 DHCPv6 global IPv6 주소 획득",
        ], accent=TEAL)

    elif t == "build":
        add_code(slide, 0.8, 1.55, 5.95, 2.05, [
            "$ cd PppRouter && make telosb",
            "$ cd TCPEcho && make telosb",
            "$ make telosb id.1 / id.2 / id.3 / id.4",
            "$ tos-bsl --telosb -c /dev/ttyUSBx ...",
        ])
        add_card(slide, 7.15, 1.55, 5.25, 2.05, "설치 결과", [
            "/dev/ttyUSB1 -> PppRouter node 1",
            "/dev/ttyUSB2 -> TCPEcho node 2",
            "/dev/ttyUSB3 -> TCPEcho node 3",
            "/dev/ttyUSB0 -> TCPEcho node 4",
        ], accent=GREEN)
        add_card(slide, 0.8, 4.1, 11.6, 1.7, "Build size", [
            "PppRouter DHCPv6 mode: ROM 48,788 bytes, RAM 9,372 bytes",
            "TCPEcho DHCPv6 mode: ROM 47,784 bytes, RAM 8,016 bytes",
            "TelosB memory budget 때문에 PppRouter의 UDP shell/route dump 제거",
        ], accent=ORANGE)

    elif t == "demo":
        add_code(slide, 0.75, 1.5, 12.0, 4.75, [
            "# 1. PPP + DHCPv6 server",
            "./start-wsn-dhcpv6.sh /dev/ttyUSB1",
            "",
            "# 2. Check leases and pick node IPv6 addresses",
            "tail -f wsn-dhcpv6.leases",
            "tail -f wsn-dhcpv6.log",
            "",
            "# 3. Commands",
            "NODE2=fd00:23:42:1::1001   # example from lease file",
            "ping6 $NODE2",
            "curl -g http://[$NODE2]/get-led",
            "curl -g http://[$NODE2]/set-led/7",
            "curl -g http://[$NODE2]/get-sensor",
        ])

    elif t == "future":
        add_card(slide, 0.75, 1.55, 3.8, 4.95, "현재 한계", [
            "현재 연결된 장비는 4 motes",
            "DHCPv6 lease를 봐야 실제 node 주소를 알 수 있음",
            "sudo 기반 PPP/dnsmasq setup은 데모 환경 의존성 존재",
        ], accent=RED)
        add_card(slide, 4.85, 1.55, 3.8, 4.95, "평가 추가", [
            "packet delivery ratio와 latency 측정",
            "hop count 변화에 따른 HTTP command RTT",
            "TCP/UDP echo 비교",
            "sensor report loss rate 측정",
        ], accent=BLUE)
        add_card(slide, 8.95, 1.55, 3.55, 4.95, "강의 기반 확장", [
            "FTSP로 get-globaltime 구현",
            "RPL MRHOF/ETX 기반 parent selection 비교",
            "QU-RPL 아이디어로 queue-aware load balancing",
            "RCRT/IFRC 관점의 congestion control 실험",
        ], accent=TEAL)

    elif t == "conclusion":
        add_label(slide, 0.9, 1.65, 3.55, 0.72, "Part-1", BLUE, 17)
        add_bullets(slide, 1.0, 2.55, 3.5, 1.75, [
            "필수 command API 구현",
            "LED, temperature, voltage, channel, TX power 반환",
        ], 17)
        add_label(slide, 4.95, 1.65, 3.55, 0.72, "Part-2", TEAL, 17)
        add_bullets(slide, 5.05, 2.55, 3.5, 1.75, [
            "DHCPv6/IPv6/6LoWPAN/RPL 통합",
            "PPP border router로 host 접근성 확보",
            "TCP/HTTP 기반 command plane",
        ], 17)
        add_label(slide, 9.0, 1.65, 3.55, 0.72, "Takeaway", ORANGE, 17)
        add_bullets(slide, 9.1, 2.55, 3.55, 1.75, [
            "강의에서 다룬 IP/DHCPv6/RPL/Transport 개념을 실제 WSN mote에 적용",
            "단순 센서 수집을 양방향 제어 가능한 WSN으로 확장",
        ], 17)

prs.save(OUT_PPTX)

notes = """# WSN Final Project 발표 대본 및 데모 체크리스트

## 1. DHCPv6 기반 IPv6/RPL Wireless Sensor Network 구현
- 이번 프로젝트의 목표는 HW2 수준의 단순 센서 수집 네트워크를 명령 기반 WSN으로 확장하는 것이다.
- node 1은 PPP border router, DHCPv6 relay, RPL root로 동작하고, node 2/3/4는 sensor node로 동작한다.
- 핵심은 WSN 내부 mote가 DHCPv6로 IPv6 주소를 자동으로 받고, 그 주소로 직접 접근 가능하게 만든 점이다.

## 2. 프로젝트 요구사항과 구현 범위
- PDF 요구사항은 base station 1개와 sensor nodes, 그리고 command 지원이다.
- 필수 명령은 get-led, set-led, get-voltage, get-temp, get-channel, get-tx-power이다.
- Part-2는 DHCPv6, RPL, TCP, Internet 연결 같은 advanced feature 중 선택할 수 있다.
- 현재 실험 환경에서는 4개 mote가 연결되어 있어 1 gateway + 3 sensor node 구성으로 데모한다.

## 3. 설계 목표
- 기존 HW2는 sensor node가 base station으로 데이터를 보내는 일방향 구조였다.
- 최종 프로젝트에서는 사용자가 node로 명령을 보내고 응답을 받아야 하므로 양방향 통신 구조가 필요했다.
- 강의의 IP narrow waist 개념처럼 IPv6를 공통 계층으로 두면 DHCPv6 주소 할당, command, TCP echo, UDP report 같은 기능을 같은 네트워크 위에서 구성할 수 있다.

## 4. 전체 시스템 구조
- PC는 PPP link를 통해 WSN에 들어가고 fd00:23:42:1::100 주소를 사용한다.
- PC의 dnsmasq가 DHCPv6 server 역할을 하며 fd00:23:42:1::/64 pool에서 주소를 배정한다.
- node 1은 PppRouter이며 PPP와 802.15.4 radio 사이에서 IPv6 packet과 DHCPv6 traffic을 forwarding/relay한다.
- node 2/3/4는 TCPEcho application을 실행하며 DHCPv6 lease로 받은 주소를 사용한다.

## 5. 프로토콜 스택과 강의 내용 연결
- IP 강의에서 배운 DHCPv6와 6LoWPAN adaptation을 사용해 작은 802.15.4 frame 위에 IPv6 packet을 실었다.
- Routing 강의에서 배운 RPL을 사용해 DODAG root와 DAO 기반 route를 구성했다.
- Transport 강의와 연결해 TCP echo와 HTTP command를 구현했고, UDP는 echo와 sensor report에 사용했다.

## 6. IPv6 주소가 없을 때 DHCPv6는 어떻게 시작되는가
- 여기서 중요한 점은 DHCPv6가 TinyOS Active Message로 따로 주소를 나눠주는 구조가 아니라는 것이다.
- sensor node는 아직 global IPv6 주소는 없지만, IEEE 802.15.4 EUI-64 기반 link-local 주소는 사용할 수 있다.
- Dhcp6ClientP는 UDP 546번 포트에 bind하고, SOLICIT 메시지를 ff02::1:2 all DHCP relay/server multicast 주소의 UDP 547번 포트로 보낸다.
- client id는 mote의 EUI-64로 만든 DUID-LL이고, DHCPv6 transaction id로 응답이 자기 요청에 대한 것인지 확인한다.
- REPLY를 받으면 IA_NA/IAADDR option에서 주소를 꺼내 SetIPAddress.setAddress()로 BLIP IPv6 stack에 global address를 등록한다.
- 따라서 bootstrap 순서는 link-local IPv6 -> DHCPv6 multicast/relay -> global IPv6 address -> application traffic이다.

## 7. DHCPv6 주소 할당 메시지 흐름
- sensor node는 SOLICIT 또는 REQUEST를 UDP 546에서 UDP 547로 보낸다.
- node1의 Dhcp6RelayP는 UDP 547에 bind되어 있고, global address가 valid해진 뒤 relay agent로 동작한다.
- node1은 sensor node의 DHCPv6 payload를 RELAY_FORW 메시지로 감싸 PC 쪽 DHCPv6 server로 전달한다.
- PC의 dnsmasq는 fd00:23:42:1::1000부터 fd00:23:42:1::1fff 범위에서 lease를 선택한다.
- server 응답은 RELAY_REPLY로 node1에 돌아오고, node1은 peer_addr을 보고 원래 sensor node에게 ADVERTISE 또는 REPLY를 전달한다.
- sensor node가 REPLY를 처리하면 그때부터 lease file에 보이는 IPv6 주소로 ping6, curl, TCP echo가 가능하다.

## 8. RPL 경로와 TCP/HTTP 동작
- DHCPv6는 주소를 주는 기능이고, RPL은 그 주소까지 packet이 갈 수 있는 경로를 만드는 기능이다.
- node1은 PppRouter에서 RootControl.setRoot()를 호출해 RPL DODAG root가 된다.
- sensor node들은 DIO를 듣고 parent를 선택하며, OF0와 storing mode 설정에 따라 default route와 downward route를 구성한다.
- PC에서 curl을 실행하면 PC -> ppp0 -> node1 -> RPL radio path 순서로 IPv6/TCP packet이 전달된다.
- sensor node에서는 TcpSocketC가 port 80 HTTP server와 port 7 TCP echo를 제공한다.
- 즉 TCP도 Active Message가 아니라 IPv6 위에서 동작하며, DHCPv6로 얻은 주소와 RPL route가 준비된 뒤 정상적으로 연결된다.

## 9. Part-1 Command 지원
- HTTP GET API로 명령을 단순화했다.
- 예를 들어 /get-led는 LED bit 상태를 decimal로 반환하고, /set-led/7은 모든 LED를 켠다.
- /get-temp, /get-voltage, /get-channel, /get-tx-power는 프로젝트 필수 명령에 대응한다.

## 10. 센서 데이터 수집 경로
- sensor node는 10초마다 temperature, humidity, voltage를 읽는다.
- 읽은 값은 cached state로 저장되어 HTTP 요청에 응답할 수 있다.
- 동시에 host의 UDP 7777 포트로 report를 보내도록 구성할 수 있으므로 host에서 주기적 sensor log를 받을 수 있다.

## 11. Part-2 Advanced Feature
- 첫 번째 advanced feature는 DHCPv6 integration이다. node들이 고정 prefix가 아니라 lease 기반 IPv6 주소를 받는다.
- 두 번째는 RPL integration이다. node 1이 root가 되고 sensor node가 DIO/DAO를 통해 route를 만든다.
- 세 번째는 TCP support이다. 각 node는 TCP echo와 HTTP command server를 제공한다.

## 12. 구현 상세
- PppRouter는 PppDaemonC, PppIpv6C, IPForwardingEngineP, RPLRoutingC, Dhcp6C를 연결한다.
- TCPEcho는 UdpSocketC, TcpSocketC, SensirionSht11C, VoltageC를 사용한다.
- TCPEcho는 StaticIPAddressTosIdC 대신 Dhcp6C를 사용해 DHCPv6로 global IPv6 주소를 획득한다.

## 13. 빌드 및 설치 결과
- PppRouter와 TCPEcho 모두 telosb target으로 빌드가 성공했다.
- 설치 매핑은 /dev/ttyUSB1 node 1, /dev/ttyUSB2 node 2, /dev/ttyUSB3 node 3, /dev/ttyUSB0 node 4이다.
- TelosB memory가 제한되어 PppRouter에서는 DHCPv6를 살리기 위해 UDP shell/route dump를 제거했다.

## 14. 데모 시나리오
1. PPP link와 DHCPv6 server 생성
```bash
./start-wsn-dhcpv6.sh /dev/ttyUSB1
```
2. lease 확인
```bash
tail -f wsn-dhcpv6.leases
tail -f wsn-dhcpv6.log
```
3. 명령 테스트
```bash
NODE2=fd00:23:42:1::1001   # 실제 값은 lease file에서 확인
ping6 $NODE2
curl -g http://[$NODE2]/get-led
curl -g http://[$NODE2]/set-led/7
curl -g http://[$NODE2]/get-temp
curl -g http://[$NODE2]/get-voltage
curl -g http://[$NODE2]/get-sensor
```

## 15. 한계와 개선 방향
- 현재 데모는 4개 mote 기준이므로, 장비가 추가되면 node 5를 같은 방식으로 확장할 수 있다.
- DHCPv6 lease를 확인해야 실제 node 주소를 알 수 있으므로, 발표 데모에서는 lease file을 먼저 보여준다.
- get-prr/get-etx, get-route를 실제 per-node link metric 기반으로 확장하면 routing 평가가 더 강해진다.
- FTSP를 붙이면 get-globaltime을 구현할 수 있고, MRHOF/QU-RPL 방식으로 routing 성능 비교도 가능하다.

## 16. 결론
- Part-1의 필수 command를 sensor node에서 HTTP API로 지원했다.
- Part-2로 DHCPv6/IPv6/6LoWPAN/RPL/TCP/PPP를 결합해 remote host가 WSN에 접근하는 구조를 구현했다.
- 강의에서 다룬 IP, DHCPv6, RPL, transport 개념을 실제 TinyOS/TelosB mote에서 동작하는 시스템으로 연결했다.
"""

with open(OUT_NOTES, "w", encoding="utf-8") as f:
    f.write(notes)

print(OUT_PPTX)
print(OUT_NOTES)
