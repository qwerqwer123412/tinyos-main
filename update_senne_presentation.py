from pathlib import Path
import shutil

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


PPTX = Path("센네 발표용.pptx")
BACKUP = Path("센네 발표용.before-dhcp-rpl-tcp.pptx")

NAVY = RGBColor(22, 45, 78)
BLUE = RGBColor(39, 98, 158)
TEAL = RGBColor(20, 132, 126)
GREEN = RGBColor(74, 148, 84)
ORANGE = RGBColor(221, 126, 49)
RED = RGBColor(188, 67, 67)
GRAY = RGBColor(92, 99, 108)
LIGHT = RGBColor(245, 248, 251)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(24, 27, 31)


def set_font(run, size=18, bold=False, color=BLACK):
    run.font.name = "Noto Sans CJK KR"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_title(slide, number, title, subtitle):
    num = slide.shapes.add_textbox(Inches(0.45), Inches(0.35), Inches(0.55), Inches(0.35))
    p = num.text_frame.paragraphs[0]
    p.text = str(number)
    set_font(p.runs[0], 13, True, GRAY)

    tb = slide.shapes.add_textbox(Inches(1.02), Inches(0.24), Inches(11.7), Inches(0.54))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    set_font(p.runs[0], 25, True, NAVY)

    sub = slide.shapes.add_textbox(Inches(1.04), Inches(0.85), Inches(11.5), Inches(0.34))
    p = sub.text_frame.paragraphs[0]
    p.text = subtitle
    set_font(p.runs[0], 12, False, GRAY)

    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.45), Inches(1.22), Inches(12.45), Inches(0.035))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()


def add_label(slide, x, y, w, h, text, fill, size=12):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    set_font(p.runs[0], size, True, WHITE)
    return shape


def add_arrow(slide, x, y, w, h, color=GRAY):
    arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()


def add_bullets(slide, x, y, w, h, bullets, size=14):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(5)
        set_font(p.runs[0], size, False, BLACK)


def add_card(slide, x, y, w, h, title, bullets, accent=BLUE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT
    shape.line.color.rgb = RGBColor(218, 225, 232)

    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.13), Inches(w - 0.3), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    set_font(p.runs[0], 15, True, NAVY)
    add_bullets(slide, x + 0.18, y + 0.55, w - 0.34, h - 0.7, bullets, 12)


def add_dhcp_bootstrap(prs, number):
    slide = prs.slides.add_slide(prs.slide_layouts[8])
    add_title(slide, number, "DHCPv6 Bootstrap: IPv6 address가 없을 때", "Active Message가 아니라 link-local IPv6 + UDP multicast로 시작")

    add_card(slide, 0.75, 1.55, 3.8, 4.95, "Initial state", [
        "Global IPv6 address는 아직 없음",
        "하지만 mote는 EUI-64 기반 link-local IPv6를 만들 수 있음",
        "BLIP IPv6 stack 위에서 DHCPv6 client가 UDP socket 사용",
        "TinyOS Active Message로 별도 주소를 배포하는 구조가 아님",
    ], BLUE)
    add_card(slide, 4.85, 1.55, 3.8, 4.95, "Client message", [
        "Dhcp6ClientP: UDP 546 bind",
        "SOLICIT -> ff02::1:2, UDP 547",
        "Client ID는 mote EUI-64 기반 DUID-LL",
        "Transaction ID로 ADVERTISE/REPLY matching",
    ], TEAL)
    add_card(slide, 8.95, 1.55, 3.55, 4.95, "After REPLY", [
        "IA_NA / IAADDR option에서 주소 추출",
        "SetIPAddress.setAddress()로 global IPv6 등록",
        "이후 ping6, TCP/HTTP, UDP report 가능",
        "lease timeout 시 removeAddress 후 갱신",
    ], ORANGE)
    return slide


def add_dhcp_relay_flow(prs, number):
    slide = prs.slides.add_slide(prs.slide_layouts[8])
    add_title(slide, number, "DHCPv6 Address Assignment Flow", "sensor node client -> node1 relay -> PC dnsmasq server")

    add_label(slide, 0.65, 1.65, 2.0, 0.62, "node2~4\nDHCPv6 Client", GREEN, 11)
    add_arrow(slide, 2.85, 1.83, 0.75, 0.26)
    add_label(slide, 3.8, 1.65, 2.05, 0.62, "node1\nRelay Agent", TEAL, 11)
    add_arrow(slide, 6.05, 1.83, 0.75, 0.26)
    add_label(slide, 7.0, 1.65, 2.35, 0.62, "PC\ndnsmasq DHCPv6", BLUE, 11)
    add_arrow(slide, 9.55, 1.83, 0.75, 0.26)
    add_label(slide, 10.5, 1.65, 2.1, 0.62, "Lease\nfd00:.../64", ORANGE, 11)

    add_card(slide, 0.75, 2.75, 5.85, 3.55, "Forward path", [
        "1. Client sends SOLICIT / REQUEST: UDP 546 -> 547",
        "2. node1 listens on UDP 547 as DHCPv6 relay",
        "3. node1 wraps payload with RELAY_FORW",
        "4. PC dnsmasq chooses address from fd00:23:42:1::1000~1fff",
    ], TEAL)
    add_card(slide, 6.85, 2.75, 5.85, 3.55, "Return path", [
        "5. Server sends RELAY_REPLY back to node1",
        "6. node1 extracts ADVERTISE / REPLY",
        "7. peer_addr tells which sensor node should receive it",
        "8. Sensor node installs the received IPv6 address",
    ], ORANGE)
    return slide


def add_rpl_tcp(prs, number):
    slide = prs.slides.add_slide(prs.slide_layouts[8])
    add_title(slide, number, "RPL route 위에서 TCP/HTTP가 동작하는 방식", "DHCPv6는 주소, RPL은 경로, TCP는 command channel")

    add_card(slide, 0.75, 1.55, 3.8, 4.95, "RPL control plane", [
        "node1 PppRouter calls RootControl.setRoot()",
        "node1이 RPL DODAG root",
        "Sensor nodes hear DIO and select parent",
        "Storing mode + DAO로 downward route 구성",
    ], GREEN)
    add_card(slide, 4.85, 1.55, 3.8, 4.95, "TCP/HTTP path", [
        "PC curl -> node IPv6:80",
        "SYN packet enters WSN through ppp0 and node1",
        "node1 forwards IPv6 packet to radio/RPL path",
        "Sensor node TcpSocketC accepts port 80 connection",
    ], BLUE)
    add_card(slide, 8.95, 1.55, 3.55, 4.95, "Key point", [
        "TCP도 Active Message가 아님",
        "IPv6 packet 안의 transport protocol로 처리",
        "RPL route가 없으면 lease가 있어도 packet delivery 불가",
        "HTTP command는 reliable TCP channel 위에서 수행",
    ], RED)
    return slide


def move_last_slides_before(prs, count, before_index):
    sld_id_lst = prs.slides._sldIdLst
    new_ids = list(sld_id_lst)[-count:]
    for sld_id in new_ids:
        sld_id_lst.remove(sld_id)
    for offset, sld_id in enumerate(new_ids):
        sld_id_lst.insert(before_index + offset, sld_id)


def renumber_visible_slide_numbers(prs):
    for idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            text = shape.text.strip()
            if text.isdigit() and len(text) <= 2:
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                p.text = str(idx)
                set_font(p.runs[0], 13, True, GRAY)
                break


def main():
    if not BACKUP.exists():
        shutil.copy2(PPTX, BACKUP)

    prs = Presentation(PPTX)
    original_count = len(prs.slides)
    conclusion_index = original_count - 1

    add_dhcp_bootstrap(prs, conclusion_index + 1)
    add_dhcp_relay_flow(prs, conclusion_index + 2)
    add_rpl_tcp(prs, conclusion_index + 3)
    move_last_slides_before(prs, 3, conclusion_index)
    renumber_visible_slide_numbers(prs)

    prs.save(PPTX)
    print(PPTX)
    print(BACKUP)
    print(f"{original_count} -> {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
